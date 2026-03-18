#!/usr/bin/env python3
"""Rearrange PowerPoint slides by index sequence.

Builds a new presentation containing slides from the source file in the
order specified by a comma-separated index list. Indices can repeat
(duplicating slides) or be omitted (removing slides).

Usage:
    python rearrange.py source.pptx output.pptx 0,3,3,5,12

Index 0 is the first slide. Repeating an index duplicates that slide.
"""

import shutil
import sys
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation


# ---------------------------------------------------------------------------
# Slide duplication via XML
# ---------------------------------------------------------------------------

def _duplicate_slide(prs, source_slide):
    """Duplicate a slide by deep-copying its XML and relationships.

    Returns the new slide object appended at the end of the deck.
    """
    # Add a slide using the same layout as the source
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Remove all auto-generated shapes from the new slide
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ""
        if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
            sp_tree.remove(child)

    # Deep-copy every shape element from source
    src_tree = source_slide.shapes._spTree
    for child in src_tree:
        tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ""
        if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
            sp_tree.append(deepcopy(child))

    # Copy image/media relationships from source slide
    for rel_key in source_slide.part.rels:
        rel = source_slide.part.rels[rel_key]
        reltype = rel.reltype

        # Image or media relationships
        if "image" in reltype or "media" in reltype:
            # Check if this rId is used in the copied XML
            if rel_key in etree.tostring(sp_tree, encoding="unicode"):
                new_slide.part.rels.get_or_add(reltype, rel.target_part)

    return new_slide


def _delete_slide(prs, slide_index):
    """Delete a slide by its 0-based index."""
    slide_id_list = prs.slides._sldIdLst
    slide_id_entry = slide_id_list[slide_index]
    rId = slide_id_entry.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    if rId:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
    slide_id_list.remove(slide_id_entry)


# ---------------------------------------------------------------------------
# Main logic
# ---------------------------------------------------------------------------

def rearrange_slides(source_path, output_path, sequence):
    """Create output.pptx with slides from source in the given order.

    Args:
        source_path: Path to source .pptx
        output_path: Path for result .pptx
        sequence: List of 0-based slide indices (may repeat)
    """
    # Work on a copy
    shutil.copy2(source_path, output_path)
    prs = Presentation(str(output_path))

    total_slides = len(prs.slides)
    for idx in sequence:
        if idx < 0 or idx >= total_slides:
            raise ValueError(
                f"Index {idx} out of range (0–{total_slides - 1}). "
                f"Source has {total_slides} slides."
            )

    # Strategy: duplicate all needed slides at the end, then delete originals.
    # This avoids issues with slide index shifting during deletion.

    original_count = total_slides

    # Phase 1: duplicate each requested slide to the end
    for idx in sequence:
        source_slide = prs.slides[idx]
        _duplicate_slide(prs, source_slide)

    # Phase 2: delete all original slides (from back to front to keep indices stable)
    for i in range(original_count - 1, -1, -1):
        _delete_slide(prs, i)

    prs.save(str(output_path))
    return len(sequence)


def main():
    if len(sys.argv) < 4:
        print(
            "Usage: python rearrange.py source.pptx output.pptx 0,3,3,5,12",
            file=sys.stderr,
        )
        sys.exit(1)

    source_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])
    sequence_str = sys.argv[3]

    if not source_path.exists():
        print(f"Error: {source_path} not found", file=sys.stderr)
        sys.exit(1)

    try:
        sequence = [int(x.strip()) for x in sequence_str.split(",") if x.strip()]
    except ValueError:
        print(f"Error: invalid sequence '{sequence_str}'. Use comma-separated integers.", file=sys.stderr)
        sys.exit(1)

    if not sequence:
        print("Error: empty sequence", file=sys.stderr)
        sys.exit(1)

    count = rearrange_slides(source_path, output_path, sequence)
    print(f"Created {output_path} with {count} slide(s) from {source_path}")


if __name__ == "__main__":
    main()

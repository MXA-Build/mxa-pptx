#!/usr/bin/env python3
"""Add native PowerPoint charts to presentation slides.

Reads a .pptx and a chart specification JSON, adds editable charts to the
specified slides, and saves the result.

Usage:
    python chart.py presentation.pptx chart-spec.json

Chart spec can be a single object or an array of objects:
{
  "slide": 2,
  "type": "bar",
  "position": {"x": 0.75, "y": 1.8, "width": 10.0, "height": 4.2},
  "data": {
    "categories": ["Q1", "Q2", "Q3", "Q4"],
    "series": [
      {"name": "Revenue", "values": [4500, 5500, 6200, 7100]}
    ]
  },
  "style": {
    "colors": ["1F4E79", "C0392B"],
    "show_legend": false,
    "show_data_labels": true
  }
}

Supported chart types: bar, column, stacked_bar, stacked_column,
100_stacked_bar, 100_stacked_column, line, pie, doughnut, scatter,
waterfall

For waterfall charts, use this data format:
{
  "type": "waterfall",
  "data": {
    "categories": ["Start", "+Growth", "-Churn", "End"],
    "values": [100, 25, -10, 115],
    "totals": [0, 3]
  },
  "style": {
    "increase_color": "2E7D32",
    "decrease_color": "C62828",
    "total_color": "1565C0"
  }
}
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Chart type mapping
# ---------------------------------------------------------------------------

_CHART_TYPES = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    "100_stacked_bar": XL_CHART_TYPE.BAR_STACKED_100,
    "100_stacked_column": XL_CHART_TYPE.COLUMN_STACKED_100,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}


def _rgb(hex_str):
    return RGBColor.from_string(hex_str.lstrip("#"))


# ---------------------------------------------------------------------------
# Standard chart builder
# ---------------------------------------------------------------------------

def _add_standard_chart(slide, spec):
    """Add a standard category-based chart (bar, column, line, pie, etc.)."""
    chart_type_name = spec["type"]
    chart_type = _CHART_TYPES.get(chart_type_name)
    if chart_type is None:
        raise ValueError(f"Unknown chart type: {chart_type_name}")

    pos = spec.get("position", {})
    x = Inches(pos.get("x", 0.75))
    y = Inches(pos.get("y", 1.8))
    w = Inches(pos.get("width", 10.0))
    h = Inches(pos.get("height", 4.2))

    data = spec.get("data", {})
    categories = data.get("categories", [])
    series_list = data.get("series", [])

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s.get("name", "Series"), s.get("values", []))

    chart_frame = slide.shapes.add_chart(chart_type, x, y, w, h, chart_data)
    chart = chart_frame.chart

    _apply_chart_style(chart, spec.get("style", {}), chart_type_name)

    return chart_frame


def _apply_chart_style(chart, style, chart_type_name):
    """Apply MXA formatting to a chart."""
    colors = style.get("colors", [])
    show_legend = style.get("show_legend", False)
    show_data_labels = style.get("show_data_labels", True)

    # Legend
    if show_legend and chart.has_legend is not None:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
    else:
        chart.has_legend = False

    # Series colours
    plot = chart.plots[0]
    for i, series in enumerate(plot.series):
        if i < len(colors):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = _rgb(colors[i])

    # Data labels
    if show_data_labels and chart_type_name not in ("scatter",):
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.color.rgb = _rgb("333333")
        if chart_type_name in ("pie", "doughnut"):
            data_labels.number_format = "0%"
        else:
            data_labels.number_format = "#,##0"

    # Gridlines — muted, value axis only
    if hasattr(chart, "value_axis"):
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = _rgb("E0E0E0")
        va.major_gridlines.format.line.width = Pt(0.5)
        va.has_minor_gridlines = False
        va.format.line.color.rgb = _rgb("CCCCCC")
        va.tick_labels.font.size = Pt(9)
        va.tick_labels.font.color.rgb = _rgb("666666")

    if hasattr(chart, "category_axis"):
        ca = chart.category_axis
        ca.has_major_gridlines = False
        ca.format.line.color.rgb = _rgb("CCCCCC")
        ca.tick_labels.font.size = Pt(9)
        ca.tick_labels.font.color.rgb = _rgb("666666")


# ---------------------------------------------------------------------------
# Waterfall chart (stacked bar with invisible base)
# ---------------------------------------------------------------------------

def _add_waterfall_chart(slide, spec):
    """Build a waterfall chart using stacked bar with invisible base series."""
    pos = spec.get("position", {})
    x = Inches(pos.get("x", 0.75))
    y = Inches(pos.get("y", 1.8))
    w = Inches(pos.get("width", 10.0))
    h = Inches(pos.get("height", 4.2))

    data = spec.get("data", {})
    categories = data.get("categories", [])
    values = data.get("values", [])
    totals = set(data.get("totals", []))  # Indices of total bars

    style = spec.get("style", {})
    increase_color = style.get("increase_color", "2E7D32")
    decrease_color = style.get("decrease_color", "C62828")
    total_color = style.get("total_color", "1565C0")

    # Calculate base/increase/decrease for each bar
    bases = []
    increases = []
    decreases = []
    running = 0

    for i, val in enumerate(values):
        if i in totals:
            # Total bar — from 0 to val
            bases.append(0)
            increases.append(val if val >= 0 else 0)
            decreases.append(abs(val) if val < 0 else 0)
            running = val
        elif val >= 0:
            bases.append(running)
            increases.append(val)
            decreases.append(0)
            running += val
        else:
            bases.append(running + val)
            increases.append(0)
            decreases.append(abs(val))
            running += val

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Base", bases)
    chart_data.add_series("Increase", increases)
    chart_data.add_series("Decrease", decreases)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, x, y, w, h, chart_data,
    )
    chart = chart_frame.chart
    plot = chart.plots[0]

    # Style the three series
    # Base — invisible
    base_series = plot.series[0]
    base_series.format.fill.background()  # No fill
    base_series.format.line.fill.background()  # No border

    # Increase — green
    inc_series = plot.series[1]
    inc_series.format.fill.solid()
    inc_series.format.fill.fore_color.rgb = _rgb(increase_color)

    # Decrease — red
    dec_series = plot.series[2]
    dec_series.format.fill.solid()
    dec_series.format.fill.fore_color.rgb = _rgb(decrease_color)

    # Recolour total bars using point-level formatting
    for idx in totals:
        if idx < len(categories):
            # For total bars, increase series holds the value (if positive)
            inc_point = inc_series.points[idx]
            inc_point.format.fill.solid()
            inc_point.format.fill.fore_color.rgb = _rgb(total_color)
            dec_point = dec_series.points[idx]
            dec_point.format.fill.solid()
            dec_point.format.fill.fore_color.rgb = _rgb(total_color)

    # No legend (series names are internal)
    chart.has_legend = False

    # Gridlines
    if hasattr(chart, "value_axis"):
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = _rgb("E0E0E0")
        va.major_gridlines.format.line.width = Pt(0.5)
        va.tick_labels.font.size = Pt(9)

    if hasattr(chart, "category_axis"):
        ca = chart.category_axis
        ca.has_major_gridlines = False
        ca.tick_labels.font.size = Pt(9)

    # Data labels on visible series
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(10)
    data_labels.number_format = "#,##0"
    # Hide labels for base series (values are 0-offset, not meaningful)
    base_series.has_data_labels = False

    return chart_frame


# ---------------------------------------------------------------------------
# Scatter chart
# ---------------------------------------------------------------------------

def _add_scatter_chart(slide, spec):
    """Add a scatter/bubble chart."""
    pos = spec.get("position", {})
    x = Inches(pos.get("x", 0.75))
    y = Inches(pos.get("y", 1.8))
    w = Inches(pos.get("width", 10.0))
    h = Inches(pos.get("height", 4.2))

    data = spec.get("data", {})
    series_list = data.get("series", [])

    chart_data = XyChartData()
    for s in series_list:
        s_data = chart_data.add_series(s.get("name", "Series"))
        x_vals = s.get("x_values", [])
        y_vals = s.get("y_values", [])
        for xv, yv in zip(x_vals, y_vals):
            s_data.add_data_point(xv, yv)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, x, y, w, h, chart_data,
    )
    chart = chart_frame.chart

    style = spec.get("style", {})
    colors = style.get("colors", [])
    plot = chart.plots[0]
    for i, series in enumerate(plot.series):
        if i < len(colors):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = _rgb(colors[i])

    return chart_frame


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

def add_chart(slide, spec):
    """Add a chart to a slide based on spec."""
    chart_type = spec.get("type", "column")

    if chart_type == "waterfall":
        return _add_waterfall_chart(slide, spec)
    elif chart_type == "scatter":
        return _add_scatter_chart(slide, spec)
    else:
        return _add_standard_chart(slide, spec)


def add_charts(pptx_path, specs, output_path=None):
    """Add one or more charts to a presentation.

    Args:
        pptx_path: Source .pptx path
        specs: List of chart spec dicts (or a single dict)
        output_path: If None, overwrites pptx_path

    Returns:
        Number of charts added
    """
    if isinstance(specs, dict):
        specs = [specs]

    prs = Presentation(str(pptx_path))
    total_slides = len(prs.slides)
    added = 0

    for spec in specs:
        slide_idx = spec.get("slide", 0)
        if slide_idx < 0 or slide_idx >= total_slides:
            print(f"Warning: slide {slide_idx} out of range (0–{total_slides-1}), skipping", file=sys.stderr)
            continue

        slide = prs.slides[slide_idx]
        add_chart(slide, spec)
        added += 1

    save_path = output_path or pptx_path
    prs.save(str(save_path))
    return added


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    if len(sys.argv) < 3:
        print("Usage: python chart.py presentation.pptx chart-spec.json [output.pptx]", file=sys.stderr)
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    spec_path = Path(sys.argv[2])
    output_path = Path(sys.argv[3]) if len(sys.argv) >= 4 else pptx_path

    if not pptx_path.exists():
        print(f"Error: {pptx_path} not found", file=sys.stderr)
        sys.exit(1)
    if not spec_path.exists():
        print(f"Error: {spec_path} not found", file=sys.stderr)
        sys.exit(1)

    specs = json.loads(spec_path.read_text(encoding="utf-8-sig"))
    count = add_charts(pptx_path, specs, output_path)
    print(f"Added {count} chart(s) → {output_path}")


if __name__ == "__main__":
    main()

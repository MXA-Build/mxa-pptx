#!/usr/bin/env python3
"""Extract structured text inventory from a PowerPoint presentation.

Outputs a JSON file mapping every slide and text-bearing shape to its
position, size, formatting, and overflow status.

Usage:
    python inventory.py input.pptx [output.json]

If output.json is omitted, prints to stdout.

Output format:
{
  "slide_width": 13.333,
  "slide_height": 7.5,
  "slides": {
    "slide-0": {
      "shape-0": {
        "name": "Title 1",
        "placeholder_type": "TITLE",
        "left": 0.75,
        "top": 0.30,
        "width": 11.83,
        "height": 1.00,
        "overflow_risk": false,
        "paragraphs": [
          {
            "text": "Revenue grew 23%",
            "alignment": "LEFT",
            "level": 0,
            "runs": [
              {
                "text": "Revenue grew 23%",
                "bold": true,
                "italic": false,
                "font_size": 24.0,
                "font_name": "Calibri",
                "color": "333333"
              }
            ]
          }
        ]
      }
    }
  }
}
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu


# ---------------------------------------------------------------------------
# Extraction helpers
# ---------------------------------------------------------------------------

EMU_PER_INCH = 914400


def _emu_to_inches(emu):
    """Convert EMU to inches, rounded to 2 decimal places."""
    if emu is None:
        return None
    return round(emu / EMU_PER_INCH, 2)


def _color_to_hex(color):
    """Try to extract hex colour string from a pptx colour object."""
    try:
        if color and color.rgb:
            return str(color.rgb)
    except (AttributeError, TypeError):
        pass
    try:
        if color and color.theme_color:
            return f"theme:{color.theme_color}"
    except (AttributeError, TypeError):
        pass
    return None


def _placeholder_type_name(shape):
    """Return the placeholder type name or None."""
    if shape.is_placeholder:
        try:
            return str(shape.placeholder_format.type).split("(")[0].strip()
        except Exception:
            return "UNKNOWN"
    return None


def _extract_run(run):
    """Extract formatting from a single text run."""
    font = run.font
    return {
        "text": run.text,
        "bold": bool(font.bold) if font.bold is not None else False,
        "italic": bool(font.italic) if font.italic is not None else False,
        "font_size": font.size.pt if font.size else None,
        "font_name": font.name,
        "color": _color_to_hex(font.color) if font.color else None,
    }


def _extract_paragraph(paragraph):
    """Extract a paragraph and its runs."""
    alignment = None
    if paragraph.alignment is not None:
        alignment = str(paragraph.alignment).split("(")[0].strip()

    return {
        "text": paragraph.text,
        "alignment": alignment,
        "level": paragraph.level if paragraph.level else 0,
        "runs": [_extract_run(r) for r in paragraph.runs],
    }


def _estimate_overflow(paragraphs, width_inches, height_inches):
    """Rough check: does text likely overflow the shape?

    Uses a simple heuristic — chars per line at a given font size.
    """
    if width_inches is None or height_inches is None:
        return False
    if width_inches <= 0 or height_inches <= 0:
        return False

    total_lines = 0
    line_height_inches = 0

    for para in paragraphs:
        # Find dominant font size in the paragraph
        sizes = []
        for run_data in para.get("runs", []):
            if run_data.get("font_size"):
                sizes.append(run_data["font_size"])
        font_size = max(sizes) if sizes else 14.0

        # ~7.5 characters per inch at 12pt, scale inversely
        chars_per_inch = 7.5 * (12.0 / font_size)
        chars_per_line = max(1, width_inches * chars_per_inch)
        text_len = len(para.get("text", ""))
        para_lines = max(1, (text_len + chars_per_line - 1) / chars_per_line)
        total_lines += para_lines
        line_height_inches = max(line_height_inches, font_size / 72.0 * 1.25)

    estimated_height = total_lines * line_height_inches
    # Allow 15% tolerance for internal padding
    return estimated_height > height_inches * 1.15


def _extract_shape(shape, shape_index):
    """Extract a text-bearing shape to dict. Returns None for non-text shapes."""
    if not shape.has_text_frame:
        return None

    paragraphs = [_extract_paragraph(p) for p in shape.text_frame.paragraphs]

    # Skip shapes with no actual text
    all_text = "".join(p["text"] for p in paragraphs).strip()
    if not all_text:
        return None

    width = _emu_to_inches(shape.width)
    height = _emu_to_inches(shape.height)

    result = {
        "name": shape.name,
        "placeholder_type": _placeholder_type_name(shape),
        "left": _emu_to_inches(shape.left),
        "top": _emu_to_inches(shape.top),
        "width": width,
        "height": height,
        "overflow_risk": _estimate_overflow(paragraphs, width, height),
        "paragraphs": paragraphs,
    }
    return result


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract_inventory(pptx_path):
    """Extract full inventory from a .pptx file.

    Returns a dict with slide dimensions and per-slide/per-shape data.
    """
    prs = Presentation(str(pptx_path))

    slide_w = _emu_to_inches(prs.slide_width)
    slide_h = _emu_to_inches(prs.slide_height)

    slides = {}
    for slide_idx, slide in enumerate(prs.slides):
        shapes = {}
        shape_counter = 0
        for shape in slide.shapes:
            data = _extract_shape(shape, shape_counter)
            if data is not None:
                shapes[f"shape-{shape_counter}"] = data
                shape_counter += 1
        if shapes:
            slides[f"slide-{slide_idx}"] = shapes

    return {
        "slide_width": slide_w,
        "slide_height": slide_h,
        "slides": slides,
    }


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    if len(sys.argv) < 2:
        print("Usage: python inventory.py input.pptx [output.json]", file=sys.stderr)
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"Error: {pptx_path} not found", file=sys.stderr)
        sys.exit(1)

    inventory = extract_inventory(pptx_path)

    output_json = json.dumps(inventory, indent=2, ensure_ascii=False)

    if len(sys.argv) >= 3:
        out_path = Path(sys.argv[2])
        out_path.write_text(output_json, encoding="utf-8")
        # Summary
        total_shapes = sum(len(s) for s in inventory["slides"].values())
        overflow_count = sum(
            1
            for slide in inventory["slides"].values()
            for shape in slide.values()
            if shape.get("overflow_risk")
        )
        print(f"Extracted {len(inventory['slides'])} slides, {total_shapes} shapes → {out_path}")
        if overflow_count:
            print(f"⚠ {overflow_count} shape(s) flagged with overflow risk")
    else:
        print(output_json)


if __name__ == "__main__":
    main()

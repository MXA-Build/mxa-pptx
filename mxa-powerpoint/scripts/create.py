#!/usr/bin/env python3
"""Create MXA-style presentations from a JSON specification.

When an MXA PowerPoint template is available, slides are built using the
template's slide layouts (inheriting the slide master's green line, logo,
and corner decoration). Content is placed using custom shapes positioned
to MXA style standards. Falls back to from-scratch generation when no
template is found.

The template slides are treated as *inspiration* -- the script does not
clone specific template slides. Instead it uses the template's layouts
(Title Slide, Numbered Title Only, Blank, etc.) and builds custom content
with shapes. This means any number of columns, any layout, can be built
while still inheriting MXA branding.

Usage:
    python create.py spec.json output.pptx
    python create.py spec.json output.pptx --template path/to/template.pptx
"""

import json
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ---------------------------------------------------------------------------
# MXA template layout indices (from slide master)
# ---------------------------------------------------------------------------

_LAYOUT_TITLE = 0            # Title Slide -- CENTER_TITLE + SUBTITLE
_LAYOUT_CONTENT = 5          # Numbered Title, Subtitle & 1-Col Content -- idx=13 + idx=14
_LAYOUT_TITLE_ONLY = 6       # Numbered Title Only -- TITLE + SLIDE_NUMBER + tab
_LAYOUT_BLANK = 16           # Blank -- SLIDE_NUMBER only
_LAYOUT_DIVIDER_DARK = 19    # Numbered Section Header Dark Green Block

_MXA_TEMPLATE_NAMES = [
    "MXA Powerpoint Template.pptx",
    "MXA Template.pptx",
    "mxa-template.pptx",
]


# ---------------------------------------------------------------------------
# MXA Style Constants (from MXA PowerPoint Style Guide)
# ---------------------------------------------------------------------------

MXA = {
    "primary":    "195B44",
    "accent":     "38B34A",
    "text":       "231F20",
    "muted":      "969696",
    "light_bg":   "F2F2F2",
    "background": "FFFFFF",
    "blue":       "248DC1",
    "purple":     "7E3794",
    "orange":     "F47E4D",
    "yellow":     "FCB53B",
}

ACCENT_ROTATION = [MXA["primary"], MXA["orange"], MXA["blue"], MXA["purple"]]
FONT = "Instrument Sans SemiBold"
BODY_FONT = "Instrument Sans"

# Positioning (inches) -- aligned with MXA template master
MARGIN_LEFT = 0.92
SLIDE_W = 13.333
CONTENT_W = SLIDE_W - MARGIN_LEFT - 0.92  # ~11.5
GREEN_LINE_Y = 1.56
CONTENT_TOP = 1.80
CONTENT_BOTTOM = 6.60
CONTENT_H = CONTENT_BOTTOM - CONTENT_TOP
FOOTER_Y = 6.92

# Font sizes (pt)
TITLE_SIZE = 24
SUBTITLE_SIZE = 20
BODY_SIZE = 16
HEADER_SIZE = 18
SOURCE_SIZE = 9
FOOTER_SIZE = 9


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rgb(hex_str):
    return RGBColor.from_string(hex_str.lstrip("#"))


def _emu(inches):
    return Emu(int(inches * 914400))


def _add_textbox(slide, x, y, w, h, text, font_size=BODY_SIZE,
                 bold=False, italic=False, color=None, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, font_name=None):
    if color is None:
        color = MXA["text"]
    if font_name is None:
        font_name = FONT
    txBox = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return txBox


def _add_bullets_box(slide, x, y, w, h, bullets, font_size=BODY_SIZE,
                     color=None, bold=False, bullet_marker=True):
    if color is None:
        color = MXA["text"]
    txBox = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)

    from lxml import etree
    _a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(6)
        p.alignment = PP_ALIGN.LEFT
        if bullet_marker:
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            pPr.set('marL', '252000')
            pPr.set('indent', '-252000')
            buClr = etree.SubElement(pPr, f'{{{_a}}}buClr')
            etree.SubElement(buClr, f'{{{_a}}}schemeClr', val='accent1')
            etree.SubElement(pPr, f'{{{_a}}}buSzPct', val="100000")
            etree.SubElement(pPr, f'{{{_a}}}buFont',
                             typeface='Wingdings',
                             panose='05000000000000000000',
                             pitchFamily='2', charset='2')
            etree.SubElement(pPr, f'{{{_a}}}buChar', char='\u00A7')
        run = p.add_run()
        run.text = bullet_text
        run.font.name = BODY_FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
    return txBox


def _add_rect(slide, x, y, w, h, fill_color, text="", font_size=HEADER_SIZE,
              font_color="FFFFFF", bold=True, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, _emu(x), _emu(y), _emu(w), _emu(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(91440)
        tf.margin_right = Emu(91440)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(font_color)
    return shape


def _ensure_footer_placeholder(slide):
    """Clone the SLIDE_NUMBER placeholder from the layout if missing on the slide."""
    # Check if it already exists
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 4:
            return
    # Clone from layout
    layout = slide.slide_layout
    for shape in layout.placeholders:
        if shape.placeholder_format.idx == 4:
            cloned = deepcopy(shape._element)
            slide.shapes._spTree.append(cloned)
            return


def _fill_footer(slide, title, date, page_num):
    _ensure_footer_placeholder(slide)
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for shape in slide.placeholders:
        if "slide number" in shape.name.lower() or shape.placeholder_format.idx == 4:
            if shape.has_text_frame:
                tf = shape.text_frame
                for p in tf.paragraphs:
                    for fld in p._p.findall("a:fld", nsmap):
                        p._p.remove(fld)
                    p.clear()
                if tf.paragraphs:
                    run = tf.paragraphs[0].add_run()
                    run.text = f" {title}    |   {date}    |    {page_num}"
                    run.font.name = FONT
                    run.font.size = Pt(FOOTER_SIZE)
                    run.font.color.rgb = _rgb(MXA["muted"])
            return


def _set_placeholder_text(slide, ph_type, text):
    from pptx.enum.shapes import PP_PLACEHOLDER
    ph_map = {
        "TITLE": PP_PLACEHOLDER.TITLE,
        "CENTER_TITLE": PP_PLACEHOLDER.CENTER_TITLE,
        "SUBTITLE": PP_PLACEHOLDER.SUBTITLE,
        "BODY": PP_PLACEHOLDER.BODY,
    }
    target = ph_map.get(ph_type)
    if target is None:
        return
    for shape in slide.placeholders:
        if shape.placeholder_format.type == target:
            if shape.has_text_frame:
                tf = shape.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    first_run = tf.paragraphs[0].runs[0]
                    while len(tf.paragraphs) > 1:
                        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
                    while len(tf.paragraphs[0].runs) > 1:
                        tf.paragraphs[0].runs[-1]._r.getparent().remove(tf.paragraphs[0].runs[-1]._r)
                    first_run.text = text
                else:
                    for p in tf.paragraphs:
                        p.clear()
                    if tf.paragraphs:
                        run = tf.paragraphs[0].add_run()
                        run.text = text
            return


def _clear_tab(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2:  # BODY
            if shape.left is not None and shape.left < _emu(0.9):
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        p.clear()
                return


def _ensure_placeholder(slide, idx):
    """Find placeholder by idx on the slide, cloning from layout if needed."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    layout = slide.slide_layout
    for shape in layout.placeholders:
        if shape.placeholder_format.idx == idx:
            cloned = deepcopy(shape._element)
            slide.shapes._spTree.append(cloned)
            for s in slide.placeholders:
                if s.placeholder_format.idx == idx:
                    return s
    return None


def _set_placeholder_subtitle(slide, text):
    """Write subtitle into placeholder idx=13 (20pt, bold, SemiBold, 231F20)."""
    ph = _ensure_placeholder(slide, 13)
    if ph is None:
        return
    tf = ph.text_frame
    for p in tf.paragraphs:
        p.clear()
    if tf.paragraphs:
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(SUBTITLE_SIZE)
        run.font.bold = True
        run.font.color.rgb = _rgb(MXA["text"])


def _set_placeholder_bullets(slide, bullets):
    """Write bullets into placeholder idx=14 with Wingdings markers."""
    ph = _ensure_placeholder(slide, 14)
    if ph is None:
        return
    tf = ph.text_frame
    tf.word_wrap = True

    from lxml import etree
    _a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Remove all existing <a:p> elements
    txBody = tf._txBody
    for existing_p in list(txBody.findall(f'{{{_a}}}p')):
        txBody.remove(existing_p)

    for bullet_text in bullets:
        p_elem = etree.SubElement(txBody, f'{{{_a}}}p')

        # <a:pPr marL="252000" indent="-252000">
        pPr = etree.SubElement(p_elem, f'{{{_a}}}pPr',
                               marL='252000', indent='-252000')

        # <a:spcBef><a:spcPts val="600"/></a:spcBef>
        spcBef = etree.SubElement(pPr, f'{{{_a}}}spcBef')
        etree.SubElement(spcBef, f'{{{_a}}}spcPts', val='600')

        # <a:buClr><a:schemeClr val="accent1"/></a:buClr>
        buClr = etree.SubElement(pPr, f'{{{_a}}}buClr')
        etree.SubElement(buClr, f'{{{_a}}}schemeClr', val='accent1')

        # <a:buSzPct val="100000"/>
        etree.SubElement(pPr, f'{{{_a}}}buSzPct', val='100000')

        # <a:buFont typeface="Wingdings" .../>
        etree.SubElement(pPr, f'{{{_a}}}buFont',
                         typeface='Wingdings',
                         panose='05000000000000000000',
                         pitchFamily='2', charset='2')

        # <a:buChar char="\u00A7"/>
        etree.SubElement(pPr, f'{{{_a}}}buChar', char='\u00A7')

        # <a:r> with sz="1800", color 231F20, no explicit font (inherits)
        r = etree.SubElement(p_elem, f'{{{_a}}}r')
        rPr = etree.SubElement(r, f'{{{_a}}}rPr', lang='en-US', dirty='0',
                               sz='1800')
        solidFill = etree.SubElement(rPr, f'{{{_a}}}solidFill')
        etree.SubElement(solidFill, f'{{{_a}}}srgbClr', val='231F20')

        t = etree.SubElement(r, f'{{{_a}}}t')
        t.text = bullet_text


# ---------------------------------------------------------------------------
# Template-based slide builders
# ---------------------------------------------------------------------------

def _build_tmpl_title(prs, spec, pres_title, date, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE])
    _set_placeholder_text(slide, "CENTER_TITLE", spec.get("lead", pres_title))
    _set_placeholder_text(slide, "SUBTITLE", spec.get("date", date))
    return slide


def _build_tmpl_exec_summary(prs, spec, pres_title, date, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_CONTENT])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", "Executive Summary"))
    _clear_tab(slide)

    subtitle = spec.get("subtitle", "")
    if subtitle:
        _set_placeholder_subtitle(slide, subtitle)

    bullets = spec.get("bullets", [])
    if bullets:
        _set_placeholder_bullets(slide, bullets)

    if spec.get("source"):
        _add_textbox(slide, MARGIN_LEFT, CONTENT_BOTTOM, CONTENT_W, 0.30,
                     f"Source: {spec['source']}", font_size=SOURCE_SIZE,
                     italic=True, color=MXA["muted"])

    _fill_footer(slide, pres_title, date, page_num)
    return slide


def _build_tmpl_divider(prs, spec, pres_title, date, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_DIVIDER_DARK])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", "Section"))

    section_num = spec.get("section_number", "")
    if section_num:
        _set_placeholder_text(slide, "BODY", str(section_num).zfill(2))

    _fill_footer(slide, pres_title, date, page_num)
    return slide


def _build_tmpl_content_text(prs, spec, pres_title, date, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_CONTENT])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", ""))
    _clear_tab(slide)

    tab = spec.get("tab", "")
    if tab:
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2 and shape.left is not None and shape.left < _emu(0.9):
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        p.clear()
                    run = shape.text_frame.paragraphs[0].add_run()
                    run.text = tab

    subtitle = spec.get("subtitle", "")
    if subtitle:
        _set_placeholder_subtitle(slide, subtitle)

    bullets = spec.get("bullets", [])
    if bullets:
        _set_placeholder_bullets(slide, bullets)

    if spec.get("source"):
        _add_textbox(slide, MARGIN_LEFT, CONTENT_BOTTOM, CONTENT_W, 0.30,
                     f"Source: {spec['source']}", font_size=SOURCE_SIZE,
                     italic=True, color=MXA["muted"])

    _fill_footer(slide, pres_title, date, page_num)
    return slide


def _build_tmpl_comparison(prs, spec, pres_title, date, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_ONLY])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", ""))
    _clear_tab(slide)

    columns = spec.get("columns", [])
    if not columns:
        _fill_footer(slide, pres_title, date, page_num)
        return slide

    num_cols = len(columns)
    gap = 0.15
    col_w = (CONTENT_W - gap * (num_cols - 1)) / num_cols
    header_h = 0.45
    body_top = CONTENT_TOP + header_h + 0.05
    body_h = CONTENT_BOTTOM - body_top - 0.10

    for i, col in enumerate(columns):
        col_x = MARGIN_LEFT + i * (col_w + gap)
        accent = ACCENT_ROTATION[i % len(ACCENT_ROTATION)]

        _add_rect(slide, col_x, CONTENT_TOP, col_w, header_h,
                  accent, text=col.get("header", f"Column {i+1}"),
                  font_size=HEADER_SIZE, font_color="FFFFFF", bold=True,
                  align=PP_ALIGN.LEFT)

        _add_rect(slide, col_x, body_top, col_w, body_h, MXA["light_bg"])

        bullets = col.get("bullets", [])
        if bullets:
            _add_bullets_box(slide, col_x, body_top, col_w, body_h,
                             bullets, font_size=BODY_SIZE, color=MXA["text"],
                             bold=False, bullet_marker=False)

    _fill_footer(slide, pres_title, date, page_num)
    return slide


def _build_tmpl_next_steps(prs, spec, pres_title, date, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_ONLY])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", "Next Steps"))
    _clear_tab(slide)

    actions = spec.get("actions", [])
    if not actions:
        _fill_footer(slide, pres_title, date, page_num)
        return slide

    col_widths = [5.5, 2.5, 2.5, 1.0]
    headers_list = ["Action", "Owner", "Timeline", "Status"]
    header_x = MARGIN_LEFT

    for i, header in enumerate(headers_list):
        _add_rect(slide, header_x, CONTENT_TOP, col_widths[i] - 0.05, 0.40,
                  MXA["primary"], text=header, font_size=12,
                  font_color="FFFFFF", bold=True)
        header_x += col_widths[i]

    for row_idx, action in enumerate(actions):
        row_y = CONTENT_TOP + 0.45 + row_idx * 0.45
        row_x = MARGIN_LEFT
        row_bg = MXA["light_bg"] if row_idx % 2 == 0 else MXA["background"]
        values = [
            action.get("action", ""),
            action.get("owner", ""),
            action.get("timeline", ""),
            action.get("status", ""),
        ]
        for i, val in enumerate(values):
            _add_rect(slide, row_x, row_y, col_widths[i] - 0.05, 0.40,
                      row_bg, text=val, font_size=12,
                      font_color=MXA["text"], bold=False)
            row_x += col_widths[i]

    _fill_footer(slide, pres_title, date, page_num)
    return slide


def _build_tmpl_appendix(prs, spec, pres_title, date, page_num):
    lead = spec.get("lead", "Appendix")
    if not lead.lower().startswith("appendix"):
        lead = f"Appendix: {lead}"
    spec = {**spec, "lead": lead}
    return _build_tmpl_content_text(prs, spec, pres_title, date, page_num)


# ---------------------------------------------------------------------------
# Content Shape builders
# ---------------------------------------------------------------------------

def _build_shape_stat_row(prs, spec, pres_title, date, page_num):
    """2-4 large stat numbers with small labels."""
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_ONLY])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", ""))
    _clear_tab(slide)

    stats = spec.get("stats", [])
    if not stats:
        _fill_footer(slide, pres_title, date, page_num)
        return slide

    n = len(stats)
    gap = 0.3
    card_w = (CONTENT_W - gap * (n - 1)) / n
    # Vertically centre the stat cards in the content zone
    card_h = 2.4
    card_y = CONTENT_TOP + (CONTENT_H - card_h) / 2

    # Auto-size stat value font to prevent line wrapping.
    # Usable text width = textbox width minus internal margins (0.1" each side).
    text_w = card_w - 0.5 - 0.2
    max_val_len = max(len(str(s.get("value", ""))) for s in stats)
    # At 54pt bold, each character is roughly 0.45" wide.
    estimated_width = max_val_len * 0.45
    stat_font = 54
    if estimated_width > text_w:
        stat_font = int(54 * text_w / estimated_width)
        stat_font = max(stat_font, 28)  # floor at 28pt

    for i, stat in enumerate(stats):
        card_x = MARGIN_LEFT + i * (card_w + gap)
        accent = ACCENT_ROTATION[i % len(ACCENT_ROTATION)]

        # Card background
        _add_rect(slide, card_x, card_y, card_w, card_h, MXA["light_bg"])

        # Left accent bar
        _add_rect(slide, card_x, card_y, 0.08, card_h, accent)

        # Large number
        _add_textbox(slide, card_x + 0.25, card_y + 0.3, card_w - 0.5, 1.2,
                     str(stat.get("value", "")),
                     font_size=stat_font, bold=True, color=accent,
                     align=PP_ALIGN.LEFT, font_name=FONT)

        # Label below
        _add_textbox(slide, card_x + 0.25, card_y + 1.5, card_w - 0.5, 0.7,
                     str(stat.get("label", "")),
                     font_size=14, bold=False, color=MXA["text"],
                     align=PP_ALIGN.LEFT, font_name=BODY_FONT)

    if spec.get("source"):
        _add_textbox(slide, MARGIN_LEFT, CONTENT_BOTTOM, CONTENT_W, 0.30,
                     f"Source: {spec['source']}", font_size=SOURCE_SIZE,
                     italic=True, color=MXA["muted"])

    _fill_footer(slide, pres_title, date, page_num)
    return slide


def _build_shape_callout_stack(prs, spec, pres_title, date, page_num):
    """2-3 stacked callout boxes with left accent borders."""
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_ONLY])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", ""))
    _clear_tab(slide)

    callouts = spec.get("callouts", [])
    if not callouts:
        _fill_footer(slide, pres_title, date, page_num)
        return slide

    n = len(callouts)
    gap = 0.2
    total_h = CONTENT_H - 0.2  # leave a little breathing room
    box_h = (total_h - gap * (n - 1)) / n
    box_h = min(box_h, 1.6)  # cap height so they don't look bloated

    for i, callout in enumerate(callouts):
        box_y = CONTENT_TOP + i * (box_h + gap)
        accent = ACCENT_ROTATION[i % len(ACCENT_ROTATION)]

        # Background card
        _add_rect(slide, MARGIN_LEFT, box_y, CONTENT_W, box_h, MXA["light_bg"])

        # Left accent bar
        _add_rect(slide, MARGIN_LEFT, box_y, 0.08, box_h, accent)

        # Text
        text = callout if isinstance(callout, str) else callout.get("text", "")
        _add_textbox(slide, MARGIN_LEFT + 0.35, box_y + 0.15,
                     CONTENT_W - 0.6, box_h - 0.3,
                     text, font_size=BODY_SIZE, bold=False,
                     color=MXA["text"], font_name=BODY_FONT,
                     anchor=MSO_ANCHOR.MIDDLE)

    if spec.get("source"):
        _add_textbox(slide, MARGIN_LEFT, CONTENT_BOTTOM, CONTENT_W, 0.30,
                     f"Source: {spec['source']}", font_size=SOURCE_SIZE,
                     italic=True, color=MXA["muted"])

    _fill_footer(slide, pres_title, date, page_num)
    return slide


def _build_shape_split(prs, spec, pres_title, date, page_num):
    """Left half bullets + right half callout/evidence box."""
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_ONLY])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", ""))
    _clear_tab(slide)

    left_w = CONTENT_W * 0.58
    right_w = CONTENT_W * 0.38
    right_x = MARGIN_LEFT + left_w + CONTENT_W * 0.04

    # Subtitle on left side
    subtitle = spec.get("subtitle", "")
    bullet_top = CONTENT_TOP
    if subtitle:
        _add_textbox(slide, MARGIN_LEFT, CONTENT_TOP, left_w, 0.4,
                     subtitle, font_size=SUBTITLE_SIZE, bold=True,
                     color=MXA["text"], font_name=FONT)
        bullet_top = CONTENT_TOP + 0.5

    # Left: bullets
    bullets = spec.get("bullets", [])
    if bullets:
        _add_bullets_box(slide, MARGIN_LEFT, bullet_top,
                         left_w, CONTENT_BOTTOM - bullet_top - 0.1,
                         bullets, font_size=BODY_SIZE, color=MXA["text"])

    # Right: callout evidence box
    callout = spec.get("callout", "")
    if callout:
        box_h = CONTENT_BOTTOM - CONTENT_TOP - 0.1
        _add_rect(slide, right_x, CONTENT_TOP, right_w, box_h, MXA["light_bg"])
        _add_rect(slide, right_x, CONTENT_TOP, 0.08, box_h, MXA["primary"])
        _add_textbox(slide, right_x + 0.25, CONTENT_TOP + 0.25,
                     right_w - 0.5, box_h - 0.5,
                     callout, font_size=15, bold=False,
                     color=MXA["text"], font_name=BODY_FONT,
                     anchor=MSO_ANCHOR.MIDDLE)

    if spec.get("source"):
        _add_textbox(slide, MARGIN_LEFT, CONTENT_BOTTOM, CONTENT_W, 0.30,
                     f"Source: {spec['source']}", font_size=SOURCE_SIZE,
                     italic=True, color=MXA["muted"])

    _fill_footer(slide, pres_title, date, page_num)
    return slide


def _build_shape_process(prs, spec, pres_title, date, page_num):
    """3-5 numbered steps connected by arrows."""
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_ONLY])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", ""))
    _clear_tab(slide)

    steps = spec.get("steps", [])
    if not steps:
        _fill_footer(slide, pres_title, date, page_num)
        return slide

    n = len(steps)
    # Arrow connectors take space between step boxes
    arrow_w = 0.35
    gap = arrow_w + 0.1
    box_w = (CONTENT_W - gap * (n - 1)) / n
    circle_r = 0.35  # radius of numbered circle
    circle_d = circle_r * 2
    step_top = CONTENT_TOP + 0.2
    desc_top = step_top + circle_d + 0.6  # below circle + title

    for i, step in enumerate(steps):
        box_x = MARGIN_LEFT + i * (box_w + gap)
        accent = ACCENT_ROTATION[i % len(ACCENT_ROTATION)]

        # Numbered circle
        cx = box_x + box_w / 2 - circle_r
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, _emu(cx), _emu(step_top),
            _emu(circle_d), _emu(circle_d)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(accent)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.name = FONT
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = _rgb("FFFFFF")

        # Step title below circle
        _add_textbox(slide, box_x, step_top + circle_d + 0.15, box_w, 0.4,
                     step.get("title", f"Step {i+1}"),
                     font_size=16, bold=True, color=accent,
                     align=PP_ALIGN.CENTER, font_name=FONT)

        # Description below title
        _add_textbox(slide, box_x, desc_top, box_w, CONTENT_BOTTOM - desc_top - 0.2,
                     step.get("description", ""),
                     font_size=14, bold=False, color=MXA["text"],
                     align=PP_ALIGN.CENTER, font_name=BODY_FONT)

        # Arrow connector (after each box except the last)
        if i < n - 1:
            arrow_x = box_x + box_w + 0.05
            arrow_y = step_top + circle_r  # middle of circle
            connector = slide.shapes.add_connector(
                1,  # straight connector
                _emu(arrow_x), _emu(arrow_y),
                _emu(arrow_x + arrow_w), _emu(arrow_y)
            )
            connector.line.color.rgb = _rgb(MXA["muted"])
            connector.line.width = Pt(2)

    if spec.get("source"):
        _add_textbox(slide, MARGIN_LEFT, CONTENT_BOTTOM, CONTENT_W, 0.30,
                     f"Source: {spec['source']}", font_size=SOURCE_SIZE,
                     italic=True, color=MXA["muted"])

    _fill_footer(slide, pres_title, date, page_num)
    return slide


def _build_shape_icon_cards(prs, spec, pres_title, date, page_num):
    """3-4 cards with coloured circle icon + header + description."""
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_ONLY])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", ""))
    _clear_tab(slide)

    cards = spec.get("cards", [])
    if not cards:
        _fill_footer(slide, pres_title, date, page_num)
        return slide

    n = len(cards)
    gap = 0.25
    card_w = (CONTENT_W - gap * (n - 1)) / n
    card_h = CONTENT_H - 0.2
    circle_d = 0.7
    icon_y = CONTENT_TOP + 0.3

    for i, card in enumerate(cards):
        card_x = MARGIN_LEFT + i * (card_w + gap)
        accent = ACCENT_ROTATION[i % len(ACCENT_ROTATION)]

        # Card background
        _add_rect(slide, card_x, CONTENT_TOP, card_w, card_h, MXA["light_bg"])

        # Icon circle
        cx = card_x + card_w / 2 - circle_d / 2
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, _emu(cx), _emu(icon_y),
            _emu(circle_d), _emu(circle_d)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(accent)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        # Use the icon letter (A, B, C) or first letter of header
        icon_text = card.get("icon", card.get("header", "?")[0])
        run.text = icon_text
        run.font.name = FONT
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = _rgb("FFFFFF")

        # Header below icon
        header_y = icon_y + circle_d + 0.2
        _add_textbox(slide, card_x + 0.15, header_y, card_w - 0.3, 0.4,
                     card.get("header", ""),
                     font_size=16, bold=True, color=accent,
                     align=PP_ALIGN.CENTER, font_name=FONT)

        # Description
        desc_y = header_y + 0.5
        _add_textbox(slide, card_x + 0.15, desc_y,
                     card_w - 0.3, CONTENT_TOP + card_h - desc_y - 0.2,
                     card.get("text", ""),
                     font_size=14, bold=False, color=MXA["text"],
                     align=PP_ALIGN.CENTER, font_name=BODY_FONT)

    if spec.get("source"):
        _add_textbox(slide, MARGIN_LEFT, CONTENT_BOTTOM, CONTENT_W, 0.30,
                     f"Source: {spec['source']}", font_size=SOURCE_SIZE,
                     italic=True, color=MXA["muted"])

    _fill_footer(slide, pres_title, date, page_num)
    return slide


def _build_shape_big_quote(prs, spec, pres_title, date, page_num):
    """Large quotation mark + quote text + attribution."""
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_ONLY])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", ""))
    _clear_tab(slide)

    quote = spec.get("quote", "")
    attribution = spec.get("attribution", "")

    if not quote:
        _fill_footer(slide, pres_title, date, page_num)
        return slide

    # Large opening quote mark
    _add_textbox(slide, MARGIN_LEFT + 0.5, CONTENT_TOP + 0.1, 1.5, 1.2,
                 "\u201C", font_size=96, bold=True, color=MXA["primary"],
                 align=PP_ALIGN.LEFT, font_name=FONT)

    # Quote text (indented, large italic)
    quote_x = MARGIN_LEFT + 1.2
    quote_w = CONTENT_W - 2.0
    _add_textbox(slide, quote_x, CONTENT_TOP + 0.9, quote_w, 2.8,
                 quote, font_size=24, bold=False, italic=True,
                 color=MXA["text"], align=PP_ALIGN.LEFT, font_name=BODY_FONT,
                 anchor=MSO_ANCHOR.TOP)

    # Attribution
    if attribution:
        _add_textbox(slide, quote_x, CONTENT_BOTTOM - 1.2, quote_w, 0.5,
                     f"\u2014 {attribution}",
                     font_size=14, bold=True, color=MXA["primary"],
                     align=PP_ALIGN.LEFT, font_name=FONT)

    if spec.get("source"):
        _add_textbox(slide, MARGIN_LEFT, CONTENT_BOTTOM, CONTENT_W, 0.30,
                     f"Source: {spec['source']}", font_size=SOURCE_SIZE,
                     italic=True, color=MXA["muted"])

    _fill_footer(slide, pres_title, date, page_num)
    return slide


def _build_shape_matrix(prs, spec, pres_title, date, page_num):
    """2x2 quadrant grid with axis labels."""
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_ONLY])
    _set_placeholder_text(slide, "TITLE", spec.get("lead", ""))
    _clear_tab(slide)

    quadrants = spec.get("quadrants", [])
    x_axis = spec.get("x_axis", "")
    y_axis = spec.get("y_axis", "")

    # Grid dimensions
    axis_label_w = 0.5
    grid_x = MARGIN_LEFT + axis_label_w + 0.1
    grid_w = CONTENT_W - axis_label_w - 0.3
    grid_y = CONTENT_TOP + 0.1
    grid_h = CONTENT_H - 0.6  # room for x-axis label below
    cell_w = grid_w / 2 - 0.05
    cell_h = grid_h / 2 - 0.05
    quad_gap = 0.1

    # Quadrant fills — top-left, top-right, bottom-left, bottom-right
    positions = [
        (grid_x, grid_y),
        (grid_x + cell_w + quad_gap, grid_y),
        (grid_x, grid_y + cell_h + quad_gap),
        (grid_x + cell_w + quad_gap, grid_y + cell_h + quad_gap),
    ]

    for i, (qx, qy) in enumerate(positions):
        accent = ACCENT_ROTATION[i % len(ACCENT_ROTATION)]
        _add_rect(slide, qx, qy, cell_w, cell_h, MXA["light_bg"])
        # Thin top accent bar
        _add_rect(slide, qx, qy, cell_w, 0.06, accent)

        if i < len(quadrants):
            q = quadrants[i]
            label = q if isinstance(q, str) else q.get("label", "")
            _add_textbox(slide, qx + 0.15, qy + 0.15, cell_w - 0.3, 0.35,
                         label, font_size=14, bold=True, color=accent,
                         font_name=FONT)
            items = [] if isinstance(q, str) else q.get("items", [])
            if items:
                items_text = "\n".join(f"\u2022 {item}" for item in items)
                _add_textbox(slide, qx + 0.15, qy + 0.55,
                             cell_w - 0.3, cell_h - 0.7,
                             items_text, font_size=12, bold=False,
                             color=MXA["text"], font_name=BODY_FONT)

    # Y-axis label (rotated text not easy in python-pptx, use vertical text)
    if y_axis:
        _add_textbox(slide, MARGIN_LEFT, grid_y + grid_h / 2 - 0.5,
                     axis_label_w, 1.0, y_axis,
                     font_size=11, bold=True, color=MXA["muted"],
                     align=PP_ALIGN.CENTER, font_name=FONT)

    # X-axis label
    if x_axis:
        _add_textbox(slide, grid_x, grid_y + grid_h + 0.05,
                     grid_w, 0.35, x_axis,
                     font_size=11, bold=True, color=MXA["muted"],
                     align=PP_ALIGN.CENTER, font_name=FONT)

    if spec.get("source"):
        _add_textbox(slide, MARGIN_LEFT, CONTENT_BOTTOM, CONTENT_W, 0.30,
                     f"Source: {spec['source']}", font_size=SOURCE_SIZE,
                     italic=True, color=MXA["muted"])

    _fill_footer(slide, pres_title, date, page_num)
    return slide


# ---------------------------------------------------------------------------
# Shape dispatch — maps shape names to builder functions
# ---------------------------------------------------------------------------

_SHAPE_BUILDERS = {
    "bullets":        None,  # handled by archetype's default builder
    "stat-row":       _build_shape_stat_row,
    "n-column":       _build_tmpl_comparison,
    "callout-stack":  _build_shape_callout_stack,
    "split":          _build_shape_split,
    "process":        _build_shape_process,
    "icon-cards":     _build_shape_icon_cards,
    "big-quote":      _build_shape_big_quote,
    "matrix":         _build_shape_matrix,
}

# Default shape for each archetype (when "shape" is not specified)
_ARCHETYPE_DEFAULT_SHAPE = {
    "comparison":    "n-column",
    "two-column":    "n-column",
    "three-column":  "n-column",
    "four-column":   "n-column",
}


# ---------------------------------------------------------------------------
# Builder dispatch (shape-aware)
# ---------------------------------------------------------------------------

# Archetype builders — used when no shape override or shape is "bullets"
_TEMPLATE_BUILDERS = {
    "title": _build_tmpl_title,
    "cover": _build_tmpl_title,
    "exec-summary": _build_tmpl_exec_summary,
    "executive-summary": _build_tmpl_exec_summary,
    "divider": _build_tmpl_divider,
    "bumper": _build_tmpl_divider,
    "content-text": _build_tmpl_content_text,
    "text": _build_tmpl_content_text,
    "content-chart": _build_tmpl_content_text,
    "chart": _build_tmpl_content_text,
    "comparison": _build_tmpl_comparison,
    "two-column": _build_tmpl_comparison,
    "three-column": _build_tmpl_comparison,
    "four-column": _build_tmpl_comparison,
    "next-steps": _build_tmpl_next_steps,
    "appendix": _build_tmpl_appendix,
}


def _resolve_builder(slide_spec):
    """Pick the right builder: shape override wins, then archetype default."""
    archetype = slide_spec.get("archetype", "content-text")
    shape = slide_spec.get("shape", "")

    # If no explicit shape, check if archetype has a default shape
    if not shape:
        shape = _ARCHETYPE_DEFAULT_SHAPE.get(archetype, "")

    # If we have a shape and it has a dedicated builder, use it
    if shape and shape != "bullets":
        builder = _SHAPE_BUILDERS.get(shape)
        if builder is not None:
            return builder

    # Fall back to archetype builder
    return _TEMPLATE_BUILDERS.get(archetype, _build_tmpl_content_text)


# ---------------------------------------------------------------------------
# From-scratch builders (fallback)
# ---------------------------------------------------------------------------

_SCRATCH_W = Inches(13.333)
_SCRATCH_H = Inches(7.5)
_SCRATCH_MARGIN = 0.75
_SCRATCH_CW = 11.83


def _scratch_lead(slide, text):
    return _add_textbox(slide, _SCRATCH_MARGIN, 0.30, _SCRATCH_CW, 1.0,
                        text, font_size=24, bold=True, color=MXA["accent"])


def _scratch_line(slide, x, y, w, color=None):
    if color is None:
        color = MXA["accent"]
    c = slide.shapes.add_connector(1, _emu(x), _emu(y), _emu(x + w), _emu(y))
    c.line.color.rgb = _rgb(color)
    c.line.width = Pt(1.5)
    return c


def _scratch_title(prs, spec, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, 13.333, 0.15, MXA["accent"])
    _add_textbox(slide, _SCRATCH_MARGIN, 2.2, _SCRATCH_CW, 1.5,
                 spec.get("lead", "Untitled"), font_size=36, bold=True,
                 color=MXA["primary"], align=PP_ALIGN.CENTER)
    _scratch_line(slide, 4.0, 3.8, 5.33)
    if spec.get("date"):
        _add_textbox(slide, _SCRATCH_MARGIN, 4.3, _SCRATCH_CW, 0.5,
                     spec["date"], font_size=14, color=MXA["muted"],
                     align=PP_ALIGN.CENTER)


def _scratch_content(prs, spec, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _scratch_lead(slide, spec.get("lead", ""))
    _scratch_line(slide, _SCRATCH_MARGIN, 1.40, _SCRATCH_CW)
    bullets = spec.get("bullets", [])
    if bullets:
        _add_bullets_box(slide, _SCRATCH_MARGIN, 1.80, _SCRATCH_CW, 4.5,
                         bullets, font_size=BODY_SIZE, color=MXA["text"])
    _add_textbox(slide, 12.0, 6.60, 0.8, 0.5, str(page_num),
                 font_size=9, color=MXA["muted"], align=PP_ALIGN.RIGHT)


def _scratch_comparison(prs, spec, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _scratch_lead(slide, spec.get("lead", ""))
    _scratch_line(slide, _SCRATCH_MARGIN, 1.40, _SCRATCH_CW)
    columns = spec.get("columns", [])
    if not columns:
        return
    n = len(columns)
    gap = 0.15
    cw = (_SCRATCH_CW - gap * (n - 1)) / n
    for i, col in enumerate(columns):
        cx = _SCRATCH_MARGIN + i * (cw + gap)
        accent = ACCENT_ROTATION[i % len(ACCENT_ROTATION)]
        _add_rect(slide, cx, 1.80, cw, 0.5, accent,
                  text=col.get("header", ""), font_size=HEADER_SIZE,
                  font_color="FFFFFF", bold=True)
        buls = col.get("bullets", [])
        if buls:
            _add_rect(slide, cx, 2.35, cw, 3.7, MXA["light_bg"])
            _add_bullets_box(slide, cx, 2.35, cw, 3.7,
                             buls, font_size=BODY_SIZE, color=MXA["text"])
    _add_textbox(slide, 12.0, 6.60, 0.8, 0.5, str(page_num),
                 font_size=9, color=MXA["muted"], align=PP_ALIGN.RIGHT)


_SCRATCH_BUILDERS = {
    "title": _scratch_title,
    "cover": _scratch_title,
    "exec-summary": _scratch_content,
    "executive-summary": _scratch_content,
    "content-text": _scratch_content,
    "text": _scratch_content,
    "comparison": _scratch_comparison,
    "next-steps": _scratch_content,
    "appendix": _scratch_content,
}


# ---------------------------------------------------------------------------
# Template search
# ---------------------------------------------------------------------------

def _find_template(spec, script_dir):
    tmpl = spec.get("template")
    if tmpl:
        p = Path(tmpl)
        if p.exists():
            return p
        p = Path.cwd() / tmpl
        if p.exists():
            return p
    for d in [Path.cwd(), script_dir.parent, script_dir.parent.parent]:
        for name in _MXA_TEMPLATE_NAMES:
            candidate = d / name
            if candidate.exists():
                return candidate
    return None


# ---------------------------------------------------------------------------
# Adjacency check
# ---------------------------------------------------------------------------

_NON_CONTENT_ARCHETYPES = {"title", "cover", "divider", "bumper", "back-cover"}


def _check_adjacency(slides_spec):
    """Warn if consecutive content slides share the same shape."""
    prev_shape = None
    for slide_spec in slides_spec:
        archetype = slide_spec.get("archetype", "content-text")
        if archetype in _NON_CONTENT_ARCHETYPES:
            prev_shape = None
            continue
        shape = slide_spec.get("shape", "")
        if not shape:
            shape = _ARCHETYPE_DEFAULT_SHAPE.get(archetype, "bullets")
        if shape == prev_shape:
            print(f"WARNING: Consecutive slides use the same shape '{shape}' "
                  f"— consider varying the visual structure for slide: "
                  f"{slide_spec.get('lead', '')[:60]}...",
                  file=sys.stderr)
        prev_shape = shape


# ---------------------------------------------------------------------------
# Spec validation
# ---------------------------------------------------------------------------

# Required fields per shape. Each entry maps to a list of (field, expected_type)
# where expected_type is 'str' for a top-level string, or 'list' for a list
# whose items must contain the listed sub-keys.
_SHAPE_REQUIRED_FIELDS = {
    "stat-row":      {"stats": [{"value", "label"}]},
    "n-column":      {"columns": [{"header"}]},
    "callout-stack": {"callouts": "list_of_str_or_dict_with_text"},
    "split":         {"bullets": "list", "callout": "str"},
    "process":       {"steps": [{"title", "description"}]},
    "icon-cards":    {"cards": [{"header", "text"}]},
    "big-quote":     {"quote": "str"},
    "matrix":        {"quadrants": "list"},
}


def _validate_slide_spec(slide_spec, slide_index):
    """Check that required fields are present for the chosen shape.

    Prints warnings to stderr for every problem found.
    Returns True if valid, False if any required field is missing or malformed.
    """
    shape = slide_spec.get("shape", "")
    if not shape:
        archetype = slide_spec.get("archetype", "content-text")
        shape = _ARCHETYPE_DEFAULT_SHAPE.get(archetype, "")
    if not shape or shape == "bullets":
        return True  # no strict schema for bullets

    reqs = _SHAPE_REQUIRED_FIELDS.get(shape)
    if reqs is None:
        return True  # unknown shape — nothing to validate

    label = slide_spec.get('lead', f'slide {slide_index + 1}')[:60]
    ok = True

    for field, constraint in reqs.items():
        val = slide_spec.get(field)

        # --- field missing or empty ---
        if val is None or val == "" or val == []:
            print(f"SPEC ERROR (slide {slide_index + 1}): shape '{shape}' "
                  f"requires '{field}' but it is missing or empty.  "
                  f"Slide: {label}",
                  file=sys.stderr)
            ok = False
            continue

        # --- list of dicts with required sub-keys ---
        if isinstance(constraint, list) and isinstance(constraint[0], set):
            required_keys = constraint[0]
            if not isinstance(val, list):
                print(f"SPEC ERROR (slide {slide_index + 1}): '{field}' must "
                      f"be a list for shape '{shape}'.  Slide: {label}",
                      file=sys.stderr)
                ok = False
                continue
            for j, item in enumerate(val):
                if isinstance(item, dict):
                    missing = required_keys - item.keys()
                    if missing:
                        print(f"SPEC ERROR (slide {slide_index + 1}): "
                              f"'{field}[{j}]' is missing keys "
                              f"{sorted(missing)} for shape '{shape}'.  "
                              f"Slide: {label}",
                              file=sys.stderr)
                        ok = False

        # --- callout-stack special: list of str or dict-with-text ---
        elif constraint == "list_of_str_or_dict_with_text":
            if not isinstance(val, list):
                print(f"SPEC ERROR (slide {slide_index + 1}): '{field}' must "
                      f"be a list for shape '{shape}'.  Slide: {label}",
                      file=sys.stderr)
                ok = False
                continue
            for j, item in enumerate(val):
                if isinstance(item, dict) and "text" not in item:
                    print(f"SPEC ERROR (slide {slide_index + 1}): "
                          f"'{field}[{j}]' is a dict but missing 'text' key "
                          f"for shape '{shape}'.  Expected keys: 'text' "
                          f"(and optional 'accent').  Got: {sorted(item.keys())}  "
                          f"Slide: {label}",
                          file=sys.stderr)
                    ok = False

        # --- simple list ---
        elif constraint == "list":
            if not isinstance(val, list):
                print(f"SPEC ERROR (slide {slide_index + 1}): '{field}' must "
                      f"be a list for shape '{shape}'.  Slide: {label}",
                      file=sys.stderr)
                ok = False

        # --- simple string ---
        elif constraint == "str":
            if not isinstance(val, str) or not val.strip():
                print(f"SPEC ERROR (slide {slide_index + 1}): '{field}' must "
                      f"be a non-empty string for shape '{shape}'.  "
                      f"Slide: {label}",
                      file=sys.stderr)
                ok = False

    return ok


# ---------------------------------------------------------------------------
# Main creation
# ---------------------------------------------------------------------------

def create_from_template(spec, template_path, output_path):
    prs = Presentation(str(template_path))
    slides_spec = spec.get("slides", [])
    pres_title = spec.get("title", "Presentation")
    date = spec.get("date", "March 2026")

    # Adjacency check — warn if consecutive content slides share the same shape
    _check_adjacency(slides_spec)

    # Delete all existing template slides (back to front)
    while len(prs.slides) > 0:
        sld_list = prs.slides._sldIdLst
        entry = sld_list[-1]
        rId = entry.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass
        sld_list.remove(entry)

    # Validate and build each slide
    for i, slide_spec in enumerate(slides_spec):
        _validate_slide_spec(slide_spec, i)
        builder = _resolve_builder(slide_spec)
        builder(prs, slide_spec, pres_title, date, i + 1)

    prs.save(str(output_path))
    return len(slides_spec)


def create_from_scratch(spec, output_path):
    prs = Presentation()
    prs.slide_width = _SCRATCH_W
    prs.slide_height = _SCRATCH_H
    slides_spec = spec.get("slides", [])
    for i, slide_spec in enumerate(slides_spec):
        archetype = slide_spec.get("archetype", "content-text")
        builder = _SCRATCH_BUILDERS.get(archetype, _scratch_content)
        builder(prs, slide_spec, i + 1)
    prs.save(str(output_path))
    return len(slides_spec)


def create_presentation(spec, output_path, template_path=None):
    script_dir = Path(__file__).resolve().parent
    if template_path is None:
        template_path = _find_template(spec, script_dir)
    if template_path and template_path.exists():
        print(f"Using template: {template_path}", file=sys.stderr)
        return create_from_template(spec, template_path, output_path)
    print("No template found -- creating from scratch", file=sys.stderr)
    return create_from_scratch(spec, output_path)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    if len(sys.argv) < 3:
        print("Usage: python create.py spec.json output.pptx [--template path.pptx]",
              file=sys.stderr)
        sys.exit(1)

    spec_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])

    template_path = None
    if "--template" in sys.argv:
        idx = sys.argv.index("--template")
        if idx + 1 < len(sys.argv):
            template_path = Path(sys.argv[idx + 1])
            if not template_path.exists():
                print(f"Error: template {template_path} not found", file=sys.stderr)
                sys.exit(1)

    if not spec_path.exists():
        print(f"Error: {spec_path} not found", file=sys.stderr)
        sys.exit(1)

    spec = json.loads(spec_path.read_text(encoding="utf-8-sig"))
    count = create_presentation(spec, output_path, template_path)
    print(f"Created {output_path} with {count} slide(s)")


if __name__ == "__main__":
    main()
#!/usr/bin/env python3
"""Apply JSON-driven text replacements to a PowerPoint presentation.

Reads a .pptx and a replacement JSON file, applies new text content to
specified shapes while preserving layout and unmodified shapes' formatting.

Usage:
    python replace.py input.pptx replacements.json output.pptx

Replacement JSON format:
{
  "slide-0": {
    "shape-0": {
      "paragraphs": [
        {
          "text": "New title text",
          "bold": true,
          "font_size": 24,
          "font_name": "Calibri",
          "color": "333333",
          "alignment": "LEFT"
        },
        {
          "text": "Second paragraph",
          "bold": false,
          "font_size": 14
        }
      ]
    }
  }
}

Shape keys (e.g. "shape-0") must match the keys from inventory.py output.
Only shapes listed in the JSON are modified; unlisted shapes are untouched.
"""

import json
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Alignment mapping
# ---------------------------------------------------------------------------

_ALIGNMENT_MAP = {
    "LEFT": PP_ALIGN.LEFT,
    "CENTER": PP_ALIGN.CENTER,
    "CENTRE": PP_ALIGN.CENTER,
    "RIGHT": PP_ALIGN.RIGHT,
    "JUSTIFY": PP_ALIGN.JUSTIFY,
}


def _parse_alignment(value):
    if value is None:
        return None
    return _ALIGNMENT_MAP.get(str(value).upper())


def _parse_color(value):
    """Parse a 6-char hex string to RGBColor."""
    if value is None:
        return None
    value = str(value).lstrip("#")
    if len(value) == 6:
        return RGBColor.from_string(value)
    return None


# ---------------------------------------------------------------------------
# Shape matching
# ---------------------------------------------------------------------------

def _build_shape_map(slide):
    """Map shape-N keys to actual shape objects that have text frames.

    Uses the same ordering as inventory.py so keys align.
    """
    mapping = {}
    counter = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            all_text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
            if all_text:
                mapping[f"shape-{counter}"] = shape
                counter += 1
    return mapping


# ---------------------------------------------------------------------------
# Replacement logic
# ---------------------------------------------------------------------------

def _clear_text_frame(text_frame):
    """Remove all text from a text frame, keeping the first paragraph element."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.text = ""
    # Clear the first paragraph text
    if text_frame.paragraphs:
        p = text_frame.paragraphs[0]
        p.clear()


def _apply_paragraphs(text_frame, para_specs):
    """Replace all text in a text frame with content from para_specs."""
    # Clear existing content
    _clear_text_frame(text_frame)

    for i, spec in enumerate(para_specs):
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()

        # Alignment
        alignment = _parse_alignment(spec.get("alignment"))
        if alignment is not None:
            para.alignment = alignment

        # Bullet level
        if "level" in spec:
            para.level = spec["level"]

        # Add run(s)
        # If "runs" key exists, use multiple runs; otherwise single run from top-level keys
        runs_spec = spec.get("runs")
        if runs_spec:
            for run_spec in runs_spec:
                run = para.add_run()
                run.text = run_spec.get("text", "")
                _apply_run_formatting(run, run_spec)
        else:
            run = para.add_run()
            run.text = spec.get("text", "")
            _apply_run_formatting(run, spec)


def _apply_run_formatting(run, spec):
    """Apply formatting from a spec dict to a run."""
    font = run.font

    if "bold" in spec:
        font.bold = bool(spec["bold"])
    if "italic" in spec:
        font.italic = bool(spec["italic"])
    if "font_size" in spec and spec["font_size"] is not None:
        font.size = Pt(spec["font_size"])
    if "font_name" in spec and spec["font_name"] is not None:
        font.name = spec["font_name"]
    if "color" in spec:
        color = _parse_color(spec["color"])
        if color is not None:
            font.color.rgb = color


# ---------------------------------------------------------------------------
# Overflow estimation (imported concept from inventory.py)
# ---------------------------------------------------------------------------

def _check_overflow(shape, para_specs):
    """Rough overflow check after replacement."""
    if shape.width is None or shape.height is None:
        return False
    width = shape.width / 914400
    height = shape.height / 914400
    if width <= 0 or height <= 0:
        return False

    total_lines = 0
    max_line_h = 0

    for spec in para_specs:
        font_size = spec.get("font_size", 14) or 14
        text = spec.get("text", "")
        chars_per_inch = 7.5 * (12.0 / font_size)
        chars_per_line = max(1, width * chars_per_inch)
        lines = max(1, (len(text) + chars_per_line - 1) / chars_per_line)
        total_lines += lines
        max_line_h = max(max_line_h, font_size / 72.0 * 1.25)

    return total_lines * max_line_h > height * 1.15


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def apply_replacements(pptx_path, replacements, output_path):
    """Apply replacement JSON to a .pptx, save to output_path.

    Returns a list of warning strings (e.g. overflow risks).
    """
    prs = Presentation(str(pptx_path))
    warnings = []
    applied = 0

    for slide_key, shapes_spec in replacements.items():
        # Parse slide index
        if not slide_key.startswith("slide-"):
            warnings.append(f"Skipping invalid key: {slide_key}")
            continue
        try:
            slide_idx = int(slide_key.split("-", 1)[1])
        except ValueError:
            warnings.append(f"Skipping invalid key: {slide_key}")
            continue

        if slide_idx < 0 or slide_idx >= len(prs.slides):
            warnings.append(f"Slide index {slide_idx} out of range (0–{len(prs.slides)-1})")
            continue

        slide = prs.slides[slide_idx]
        shape_map = _build_shape_map(slide)

        for shape_key, shape_spec in shapes_spec.items():
            if shape_key not in shape_map:
                warnings.append(
                    f"{slide_key}/{shape_key}: shape not found. "
                    f"Available: {list(shape_map.keys())}"
                )
                continue

            shape = shape_map[shape_key]
            para_specs = shape_spec.get("paragraphs", [])
            if not para_specs:
                warnings.append(f"{slide_key}/{shape_key}: no paragraphs in spec, skipping")
                continue

            # Check overflow before applying
            if _check_overflow(shape, para_specs):
                warnings.append(
                    f"{slide_key}/{shape_key}: text may overflow shape "
                    f"({shape.width/914400:.1f}\" × {shape.height/914400:.1f}\")"
                )

            _apply_paragraphs(shape.text_frame, para_specs)
            applied += 1

    prs.save(str(output_path))
    return applied, warnings


def main():
    if len(sys.argv) < 4:
        print(
            "Usage: python replace.py input.pptx replacements.json output.pptx",
            file=sys.stderr,
        )
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    json_path = Path(sys.argv[2])
    output_path = Path(sys.argv[3])

    if not pptx_path.exists():
        print(f"Error: {pptx_path} not found", file=sys.stderr)
        sys.exit(1)
    if not json_path.exists():
        print(f"Error: {json_path} not found", file=sys.stderr)
        sys.exit(1)

    replacements = json.loads(json_path.read_text(encoding="utf-8-sig"))

    # Work on a copy so input is never modified
    shutil.copy2(pptx_path, output_path)
    applied, warnings = apply_replacements(output_path, replacements, output_path)

    print(f"Applied {applied} replacement(s) → {output_path}")
    for w in warnings:
        print(f"  ⚠ {w}")

    if warnings:
        sys.exit(2)  # Non-zero but not 1 (which is fatal error)


if __name__ == "__main__":
    main()
